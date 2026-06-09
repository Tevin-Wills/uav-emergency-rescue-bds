#!/usr/bin/env python3
"""Build Level3_Presentation.pptx — 4-slide poster-talk deck.
Design per the 6-paper-slides skill: figure slides image >= 60% area,
title >= 28pt (primary), body >= 20pt, bold key numbers, frame numbers,
per-slide speaker notes, clean white background.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, 'figures')
NAVY = RGBColor(0x1F, 0x49, 0x7B)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
DARK = RGBColor(0x22, 0x2B, 0x33)
GREY = RGBColor(0x80, 0x80, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
blank = prs.slide_layouts[6]


def add_title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.33), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = NAVY
    r.font.name = 'Calibri'
    return tb


def add_rule(slide, y=1.18):
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.33), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()


def fit_image(slide, path, max_w, max_h, cx, top):
    iw, ih = Image.open(path).size
    asp = iw / ih
    w = max_w; h = w / asp
    if h > max_h:
        h = max_h; w = h * asp
    left = cx - w / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return w * h  # area sq in


def bullets(slide, items, left, top, width, height, size=20, gap=6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, bolds) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # split on ** for bold spans
        parts = txt.split('**')
        for j, seg in enumerate(parts):
            if seg == '':
                continue
            r = p.add_run(); r.text = ('• ' if j == 0 else '') + seg if j == 0 else seg
            r.font.size = Pt(size); r.font.name = 'Calibri'
            r.font.color.rgb = DARK
            if j % 2 == 1:
                r.font.bold = True; r.font.color.rgb = NAVY
    return tb


def caption(slide, segs, top=6.82):
    """Single centered caption line; segs = list of (text, bold)."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.33), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for txt, bold in segs:
        r = p.add_run(); r.text = txt; r.font.size = Pt(18); r.font.name = 'Calibri'
        r.font.bold = bool(bold); r.font.color.rgb = ACCENT if bold else DARK
    return tb


def frame_num(slide, n):
    tb = slide.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f'{n} / 4'; r.font.size = Pt(11); r.font.color.rgb = GREY


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
# accent band
band = s.shapes.add_shape(1, Inches(0), Inches(2.45), Inches(13.333), Inches(0.10))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.7), Inches(1.0))
p = tb.text_frame.paragraphs[0]; r = p.add_run()
r.text = 'Resilient RTK Positioning for Disaster Rescue'
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.65), Inches(11.7), Inches(0.7))
p = tb.text_frame.paragraphs[0]; r = p.add_run()
r.text = 'Level 3 — RTK under a compound GNSS-degradation disaster'
r.font.size = Pt(22); r.font.italic = True; r.font.color.rgb = DARK

# hook line with bold numbers
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
for seg, bold in [('Mean error held to ', 0), ('0.968 m', 1),
                  (' under disaster conditions — ', 0), ('5x', 1),
                  (' better than total GNSS failure. Approach navigation ', 0),
                  ('PASSED', 1), (' (92.4% within 2 m).', 0)]:
    r = p.add_run(); r.text = seg; r.font.size = Pt(22)
    r.font.bold = bool(bold); r.font.color.rgb = ACCENT if bold else DARK

tb = s.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
p = tb.text_frame.paragraphs[0]; r = p.add_run()
r.text = 'UAV Emergency Rescue BDS  ·  RTK Positioning Module  ·  Tevin Wills  ·  2026-06-04'
r.font.size = Pt(14); r.font.color.rgb = GREY
frame_num(s, 1)
notes(s, "Open with the problem, not the method. We needed RTK to survive a disaster zone "
         "where GNSS degrades on multiple axes at once. Headline: even under the compound "
         "disaster, RTK held mean error under 1 metre and passed the approach target. "
         "The honest gap is precision landing, which I'll come to on slide 3.")

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Method (scenario profile figure)
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, 'The compound disaster we designed')
add_rule(s)
fit_image(s, os.path.join(FIG, 'l3_compound_scenario_profile.png'),
          max_w=12.8, max_h=5.4, cx=SW / 2, top=1.35)
caption(s, [('Three co-occurring degradations: noise ', 0), ('1.5 -> 2.75 m', 1),
            (',  quality ', 0), ('0.95 -> 0.60', 1), (',  dropouts ', 0),
            ('37.5 s / 90 s', 1), ('  across 5 phases', 0)])
frame_num(s, 2)
notes(s, "This is the input model, not a result. Real disaster zones don't fail one way at "
         "a time, so we degrade noise, correction quality and correction availability together, "
         "progressively, across the five mission phases. Peak degradation is the Search Zone "
         "over the collapse site: 2.75 m noise, 0.60 quality, and 37.5-second correction "
         "dropouts every 90 seconds.")

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Results (three-run comparison figure)
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, 'What we confirmed')
add_rule(s)
fit_image(s, os.path.join(FIG, 'l3_three_run_comparison.png'),
          max_w=12.6, max_h=4.2, cx=SW / 2, top=1.35)
bullets(s, [
    ('Graceful degradation: compound **0.968 m** vs total failure **4.776 m**  (~**5x** better)', None),
    ('Approach navigation **PASSED** — **92.4%** within 2.0 m (target 90%)', None),
    ('Precision landing **FELL SHORT** — **31.5%** within 0.3 m (target 80%)', None),
], left=0.6, top=5.8, width=12.2, height=1.5, size=18, gap=4)
frame_num(s, 3)
notes(s, "Three runs bracket the result: ideal baseline at 5 cm, total failure at nearly 5 m, "
         "and the compound disaster in between at under 1 m — about five times better than "
         "having no corrections. We passed the approach target comfortably. The honest finding "
         "is precision landing: we hold approach-grade accuracy but only reach the 0.3 m bar "
         "31.5% of the time during the disaster, because we're in RTK Float, not Fixed.")

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Takeaway & next steps
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, 'Takeaway & next steps')
add_rule(s)

tb = s.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(11.9), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = 'Confirmed'
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = GREEN
for txt in [
    'RTK degrades gracefully — usable accuracy under sustained degradation',
    'Approach viable (92.4%) + honest, dynamic uncertainty + correct viability states',
    'Data proven real — PX4 ULog and analysis CSV agree on the trajectory',
]:
    pp = tf.add_paragraph(); pp.space_after = Pt(5)
    rr = pp.add_run(); rr.text = '• ' + txt; rr.font.size = Pt(19); rr.font.color.rgb = DARK

tb = s.shapes.add_textbox(Inches(0.7), Inches(4.15), Inches(11.9), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = 'Next'
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = ACCENT
for txt in [
    'Close the landing gap — RTK Fixed re-lock / re-convergence in the landing phase',
    'Re-fly both runs with one identical mission for strict comparability',
    'Phase 3 — validate the noise model against real BeiDou data via RTKLIB',
]:
    pp = tf.add_paragraph(); pp.space_after = Pt(5)
    rr = pp.add_run(); rr.text = '• ' + txt; rr.font.size = Pt(19); rr.font.color.rgb = DARK
frame_num(s, 4)
notes(s, "To close: three things confirmed — graceful degradation, viable approach navigation "
         "with an honest uncertainty signal, and provably real flight data. The single open "
         "engineering gap is precision landing, and the next step is a Fixed-lock recovery "
         "protocol for the landing phase. Beyond that, Phase 3 grounds the whole noise model "
         "in real BeiDou measurements through RTKLIB. Thank you.")

out = os.path.join(HERE, 'Level3_Presentation.pptx')
prs.save(out)
print('saved', out, '| slides:', len(prs.slides._sldIdLst))
